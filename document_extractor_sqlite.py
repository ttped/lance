# document_extractor_sqlite.py

import hashlib
import mimetypes
import sqlite3
import json # For serializing content_units to JSON string
from datetime import datetime, timezone
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple

import PyPDF2
from pptx import Presentation

# --- Configuration ---
SQLITE_DATABASE_FILE = "engineering_docs.sqlite"
PROCESSED_DOCUMENTS_TABLE_NAME = "processed_documents"
SUPPORTED_EXTENSIONS = ['.txt', '.pdf', '.pptx']
MAX_CONFIRMATION_PRINTOUTS = 5 # Max number of files to print details for


class DocumentExtractorSQLite:
    def __init__(self, db_path: str):
        self.db_path = db_path
        self._initialize_database()

    def _initialize_database(self):
        """Creates the database and table if they don't exist."""
        try:
            with sqlite3.connect(self.db_path) as conn:
                cursor = conn.cursor()
                cursor.execute(f"""
                CREATE TABLE IF NOT EXISTS {PROCESSED_DOCUMENTS_TABLE_NAME} (
                    id INTEGER PRIMARY KEY AUTOINCREMENT,
                    filename TEXT NOT NULL,
                    file_path_original TEXT NOT NULL,
                    file_type TEXT,
                    file_hash_sha256 TEXT NOT NULL UNIQUE,
                    processed_timestamp TEXT NOT NULL,
                    total_units INTEGER,
                    content_units_json TEXT
                )
                """)
                conn.commit()
        except sqlite3.Error as e:
            print(f"Database initialization error: {e}")
            raise

    def _calculate_file_hash(self, file_path: Path) -> str:
        """Calculates the SHA256 hash of a file."""
        sha256_hash = hashlib.sha256()
        with open(file_path, "rb") as f:
            for byte_block in iter(lambda: f.read(4096), b""):
                sha256_hash.update(byte_block)
        return sha256_hash.hexdigest()

    def _extract_text_from_txt(self, file_path: Path) -> List[Dict[str, Any]]:
        """Extracts text from a TXT file."""
        content_units = []
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
                if text:
                    content_units.append({
                        "unit_number": 1, # TXT file treated as a single unit
                        "text": text,
                        "char_count": len(text)
                    })
        except Exception as e:
            print(f"An unexpected error occurred while processing TXT {file_path.name}: {e}")
        return content_units

    def _extract_text_from_pdf(self, file_path: Path) -> List[Dict[str, Any]]:
        """Extracts text from each page of a PDF file."""
        content_units = []
        try:
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for i, page in enumerate(reader.pages):
                    text = page.extract_text()
                    if text: # Only add if text was extracted
                        content_units.append({
                            "unit_number": i + 1,
                            "text": text,
                            "char_count": len(text)
                        })
        except PyPDF2.errors.PdfReadError as e:
            print(f"Error reading PDF {file_path.name}: {e}. It might be encrypted or corrupted.")
        except Exception as e:
            print(f"An unexpected error occurred while processing PDF {file_path.name}: {e}")
        return content_units

    def _extract_text_from_pptx(self, file_path: Path) -> List[Dict[str, Any]]:
        """Extracts text from each slide of a PowerPoint file."""
        content_units = []
        try:
            prs = Presentation(file_path)
            for i, slide in enumerate(prs.slides):
                slide_text_parts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                slide_text_parts.append(run.text)
                
                full_slide_text = "\n".join(slide_text_parts).strip()
                if full_slide_text: # Only add if text was extracted
                    content_units.append({
                        "unit_number": i + 1,
                        "text": full_slide_text,
                        "char_count": len(full_slide_text)
                    })
        except Exception as e:
            print(f"An error occurred while processing PPTX {file_path.name}: {e}")
        return content_units

    def process_file(self, file_path_str: str) -> Optional[Dict[str, Any]]:
        """
        Processes a single file: extracts text and stores it in SQLite
        if not already processed.
        Returns a dictionary with processing details if newly added,
        the existing row ID if skipped, or None if failed.
        """
        file_path = Path(file_path_str)
        if not file_path.exists() or not file_path.is_file():
            print(f"File not found or is not a file: {file_path_str}")
            return None # Indicates failure or file not found

        file_hash = self._calculate_file_hash(file_path)

        try:
            with sqlite3.connect(self.db_path) as conn:
                cursor = conn.cursor()
                cursor.execute(
                    f"SELECT id FROM {PROCESSED_DOCUMENTS_TABLE_NAME} WHERE file_hash_sha256 = ?",
                    (file_hash,)
                )
                existing_row = cursor.fetchone()
                if existing_row:
                    print(f"File {file_path.name} (hash: {file_hash[:8]}...) already processed. Row ID: {existing_row[0]}. Skipping.")
                    return {"status": "skipped", "row_id": existing_row[0], "filename": file_path.name}


                # Determine file type based on extension for simplicity here,
                # though mime_type could also be used.
                file_extension = file_path.suffix.lower()
                file_type = ""
                content_units = []

                if file_extension == ".pdf":
                    file_type = "pdf"
                    content_units = self._extract_text_from_pdf(file_path)
                elif file_extension == ".pptx":
                    file_type = "pptx"
                    content_units = self._extract_text_from_pptx(file_path)
                elif file_extension == ".txt":
                    file_type = "txt"
                    content_units = self._extract_text_from_txt(file_path)
                else:
                    # This check is largely redundant if we filter by SUPPORTED_EXTENSIONS upstream
                    print(f"Unsupported file extension: {file_extension} for file {file_path.name}. Skipping.")
                    return {"status": "unsupported", "filename": file_path.name}


                if not content_units:
                    print(f"No text content extracted from {file_path.name}. Skipping database insertion.")
                    return {"status": "no_content", "filename": file_path.name}

                content_units_json = json.dumps(content_units)
                timestamp_now = datetime.now(timezone.utc).isoformat()

                insert_query = f"""
                INSERT INTO {PROCESSED_DOCUMENTS_TABLE_NAME} (
                    filename, file_path_original, file_type, file_hash_sha256,
                    processed_timestamp, total_units, content_units_json
                ) VALUES (?, ?, ?, ?, ?, ?, ?)
                """
                params = (
                    file_path.name,
                    str(file_path.resolve()),
                    file_type,
                    file_hash,
                    timestamp_now,
                    len(content_units),
                    content_units_json
                )
                
                cursor.execute(insert_query, params)
                conn.commit()
                inserted_id = cursor.lastrowid
                print(f"Successfully processed and stored {file_path.name}. Row ID: {inserted_id}")
                
                return {
                    "status": "added",
                    "row_id": inserted_id,
                    "filename": file_path.name,
                    "file_type": file_type,
                    "total_units": len(content_units),
                    "content_units_for_snippet": content_units # Pass this for snippet generation
                }

        except sqlite3.Error as e:
            print(f"SQLite error for {file_path.name}: {e}")
            return {"status": "error", "filename": file_path.name, "details": str(e)}
        except Exception as e:
            print(f"An unexpected error occurred processing {file_path.name}: {e}")
            return {"status": "error", "filename": file_path.name, "details": str(e)}


def scan_directory_and_process(extractor: DocumentExtractorSQLite, target_dir_path: Path):
    """Scans a directory recursively and processes supported files."""
    if not target_dir_path.is_dir():
        print(f"Error: Target path '{target_dir_path}' is not a directory or does not exist.")
        return

    print(f"\nScanning directory: {target_dir_path} for files with extensions: {SUPPORTED_EXTENSIONS}")
    
    processed_for_confirmation_count = 0
    files_found_count = 0
    files_processed_count = 0
    files_skipped_count = 0
    files_error_count = 0

    for file_path in target_dir_path.rglob('*'): # rglob for recursive
        if file_path.is_file() and file_path.suffix.lower() in SUPPORTED_EXTENSIONS:
            files_found_count += 1
            print(f"\nFound supported file: {file_path}")
            result = extractor.process_file(str(file_path))

            if result:
                if result.get("status") == "added":
                    files_processed_count +=1
                    if processed_for_confirmation_count < MAX_CONFIRMATION_PRINTOUTS:
                        print("-" * 30)
                        print(f"CONFIRMATION FOR NEWLY ADDED FILE:")
                        print(f"  Filename:    {result['filename']}")
                        print(f"  File Type:   {result['file_type']}")
                        print(f"  Row ID:      {result['row_id']}")
                        print(f"  Total Units: {result['total_units']}")
                        
                        content_units = result.get("content_units_for_snippet", [])
                        if content_units and len(content_units) > 0:
                            first_unit_text = content_units[0].get("text", "")
                            snippet = (first_unit_text[:150] + '...') if len(first_unit_text) > 150 else first_unit_text
                            print(f"  Text Snippet (Unit 1):\n    '{snippet.replace(chr(10), chr(10) + '     ')}'")
                        else:
                            print("  Text Snippet: (No text extracted or no units)")
                        print("-" * 30)
                        processed_for_confirmation_count += 1
                elif result.get("status") == "skipped":
                    files_skipped_count +=1
                elif result.get("status") == "error" or result.get("status") == "unsupported" or result.get("status") == "no_content":
                    files_error_count +=1
            else: # Should not happen if process_file always returns a dict, but as a safeguard
                files_error_count +=1


    print("\n--- Processing Summary ---")
    print(f"Total supported files found: {files_found_count}")
    print(f"Successfully processed and added: {files_processed_count}")
    print(f"Skipped (already in DB): {files_skipped_count}")
    print(f"Errors or no content: {files_error_count}")
    print(f"Details printed for {processed_for_confirmation_count} newly added file(s).")
    print("--- End of Summary ---")


if __name__ == "__main__":
    print(f"Starting Document Extractor with SQLite backend (DB: {SQLITE_DATABASE_FILE})...")
    
    # --- TARGET DIRECTORY FOR PROCESSING ---
    # Update this path to the directory you want to scan
    directory_to_scan = Path("./dummy_files_dir") 
    # --- --- --- --- --- --- --- --- --- ---

    directory_to_scan.mkdir(exist_ok=True) # Ensure dummy_files_dir exists for the example

    try:
        extractor = DocumentExtractorSQLite(SQLITE_DATABASE_FILE)
        scan_directory_and_process(extractor, directory_to_scan)
        
        # Example of how to add more files or directories:
        # another_directory_to_scan = Path("./another_folder_with_docs")
        # if another_directory_to_scan.is_dir():
        #     scan_directory_and_process(extractor, another_directory_to_scan)

    except Exception as e:
        print(f"Failed to initialize or run the extractor: {e}")
        import traceback
        traceback.print_exc()

    print("\nDocument extraction process finished.")
